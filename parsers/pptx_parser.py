from __future__ import annotations

import base64

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import ParsedNode, ParsedTable


def parse_pptx(path: str) -> tuple[list[ParsedNode], int]:
    prs = Presentation(path)
    nodes: list[ParsedNode] = []
    image_count = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape and title_shape.text else f"第 {slide_idx} 页"
        node = ParsedNode(level=1, title=title)

        for shape in slide.shapes:
            if shape == title_shape:
                continue

            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text = (shape.text or "").strip()
                if text:
                    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
                    if len(lines) == 1:
                        node.bullets.append({"title": lines[0], "description": ""})
                    else:
                        node.raw_text = (node.raw_text + "\n" + "\n".join(lines)).strip()

            if getattr(shape, "has_table", False):
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([(cell.text or "").strip() for cell in row.cells])
                if rows:
                    node.tables.append(
                        ParsedTable(headers=rows[0], rows=rows[1:] if len(rows) > 1 else [], caption="")
                    )

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                image_blob = getattr(getattr(shape, "image", None), "blob", b"") or b""
                image_ext = (getattr(getattr(shape, "image", None), "ext", "") or "").lower()
                if image_ext in {"jpg", "jpeg"}:
                    mime = "image/jpeg"
                elif image_ext == "png":
                    mime = "image/png"
                elif image_ext == "gif":
                    mime = "image/gif"
                elif image_ext == "webp":
                    mime = "image/webp"
                elif image_ext == "bmp":
                    mime = "image/bmp"
                else:
                    mime = "image/png"
                data_uri = ""
                if image_blob:
                    b64 = base64.b64encode(image_blob).decode("ascii")
                    data_uri = f"data:{mime};base64,{b64}"
                node.images.append(
                    {"path": data_uri, "caption": f"幻灯片图片 {image_count}", "source": "pptx"}
                )

        nodes.append(node)

    return nodes, image_count
